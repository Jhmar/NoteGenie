from fastapi import UploadFile, File
from fastapi.responses import JSONResponse
import os
from src.utils.ocr_utils import extract_text_from_file, setup_tesseract
from src.utils.file_utils import save_temp_file, cleanup_temp_files

setup_tesseract()

async def extract_text(file: UploadFile = File(...)):
    """Extract text from image, PDF, PPTX, DOCX, or TXT."""
    temp_path, file_extension = save_temp_file(file)
    try:
        ext = file_extension.lower()

        if ext == '.pptx':
            text = _extract_pptx(temp_path)
            return JSONResponse({"filename": file.filename, "file_type": "PPTX",
                "extracted_text": text, "status": "success", "characters_extracted": len(text)})

        if ext == '.docx':
            text = _extract_docx(temp_path)
            return JSONResponse({"filename": file.filename, "file_type": "DOCX",
                "extracted_text": text, "status": "success", "characters_extracted": len(text)})

        if ext == '.txt':
            with open(temp_path, 'r', errors='ignore') as f:
                text = f.read()
            return JSONResponse({"filename": file.filename, "file_type": "TXT",
                "extracted_text": text, "status": "success", "characters_extracted": len(text)})

        result = extract_text_from_file(temp_path, file_extension)
        return JSONResponse({"filename": file.filename, "file_type": result["type"],
            "extracted_text": result["text"],
            "status": "success" if result["type"] in ["image", "PDF"] else "error",
            "characters_extracted": result["characters"]})

    except Exception as e:
        return JSONResponse({"filename": file.filename, "error": str(e), "status": "error"}, status_code=500)
    finally:
        cleanup_temp_files([temp_path])


async def upload_file(file: UploadFile = File(...)):
    """Simple file upload endpoint."""
    from src.utils.config import UPLOAD_DIR
    from src.utils.file_utils import save_uploaded_file
    file_path = save_uploaded_file(file, UPLOAD_DIR)
    return JSONResponse({"filename": file.filename, "size": file_path.stat().st_size, "message": "File uploaded successfully"})


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
    return "\n".join(lines).strip()


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())